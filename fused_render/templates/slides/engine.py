# /// script
# dependencies = ["python-pptx", "Pillow", "fpdf2"]
# ///
"""DOM-free slide engine for the slides template.

This module owns the two hard directions and the canonical edit semantics, with
NO knowledge of HTTP, the browser, or fused-render. It is import-safe and unit
testable on its own (`python -m fused_render.templates.slides.engine <file.pptx>`).

  pptx bytes  --parse-->  model.json  (canonical, AI-addressable)
  model.json  --build-->  pptx bytes

The **model** is the single source of truth the AI agent edits. Every slide and
every element carries a stable string id (`s_*`, `e_*`) so an agent can target
"the title on slide 2" without guessing pixel coordinates. Coordinates live in a
slide-pixel space (EMU / 9525 == 96dpi px); a standard 16:9 deck is 1280x720.

Element types kept deliberately small so the renderer and the round-trip stay
honest:
  text   - a box of paragraphs/runs (covers title/body placeholders, text boxes,
           and autoshapes, which may additionally carry a fill/outline)
  image  - a raster picture stored under the deck's media/ folder
  table  - rows x cols of cell text
Anything we cannot faithfully rebuild (charts, groups, smartart) is captured as
a text box labelled with its kind, so nothing silently vanishes from the canvas.
"""

from __future__ import annotations

import os
import uuid

EMU_PER_PX = 9525  # 914400 EMU/inch / 96 px/inch
DEFAULT_W_PX = 1280
DEFAULT_H_PX = 720
# Bump when parse output changes so cached decks auto-reparse (folded into the
# cache key in slides.py, since caches there are keyed by content hash).
ENGINE_V = 2


def _nid(prefix: str) -> str:
    return f"{prefix}_{uuid.uuid4().hex[:8]}"


def emu_to_px(v) -> float:
    if v is None:
        return 0.0
    return round(v / EMU_PER_PX, 1)


def px_to_emu(v) -> int:
    return int(round(float(v) * EMU_PER_PX))


# --------------------------------------------------------------------------- #
#  colour helpers                                                             #
# --------------------------------------------------------------------------- #
def _rgb_hex(color) -> str | None:
    """Best-effort solid RGB out of a python-pptx color/fill object.

    Theme-referenced colours have no absolute rgb; we return None and let the
    caller fall back to a sensible default rather than raising.
    """
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE

        if color is None:
            return None
        if getattr(color, "type", None) == MSO_COLOR_TYPE.RGB:
            return "#" + str(color.rgb)
    except Exception:
        pass
    return None


def _fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
        from pptx.enum.dml import MSO_FILL_TYPE

        if fill.type == MSO_FILL_TYPE.SOLID:
            return _rgb_hex(fill.fore_color)
    except Exception:
        pass
    return None


def _rpr_props(el) -> dict:
    """Pull run props (size pt, bold, italic, font, color) from an <a:defRPr>/<a:rPr>."""
    out = {}
    if el is None:
        return out
    sz = el.get("sz")
    if sz:
        try:
            out["size"] = round(int(sz) / 100.0, 1)
        except Exception:
            pass
    for attr, key in (("b", "bold"), ("i", "italic")):
        v = el.get(attr)
        if v is not None:
            out[key] = v in ("1", "true")
    from pptx.oxml.ns import qn

    latin = el.find(qn("a:latin"))
    if latin is not None:
        tf = latin.get("typeface")
        if tf and not tf.startswith("+"):  # skip theme refs (+mn-lt/+mj-lt)
            out["font"] = tf
    sf = el.find(qn("a:solidFill"))
    if sf is not None:
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is not None and srgb.get("val"):
            out["color"] = "#" + srgb.get("val")
    return out


def _ph_props(ph, level) -> dict:
    """Inherited run props from a placeholder's txBody (lstStyle lvlN, then pPr)."""
    from pptx.oxml.ns import qn

    props = {}
    try:
        tb = ph.text_frame._txBody
    except Exception:
        return props
    lst = tb.find(qn("a:lstStyle"))
    if lst is not None:
        lvl = lst.find(qn("a:lvl%dpPr" % (min(level, 8) + 1)))
        if lvl is not None:
            props.update(_rpr_props(lvl.find(qn("a:defRPr"))))
    for p in tb.findall(qn("a:p")):
        pPr = p.find(qn("a:pPr"))
        if pPr is not None:
            for k, v in _rpr_props(pPr.find(qn("a:defRPr"))).items():
                props.setdefault(k, v)
        break
    return props


def _txstyle_props(master, cat, level) -> dict:
    from pptx.oxml.ns import qn

    tx = master._element.find(qn("p:txStyles"))
    if tx is None:
        return {}
    style = tx.find(qn({"title": "p:titleStyle", "body": "p:bodyStyle"}.get(cat, "p:otherStyle")))
    if style is None:
        return {}
    lvl = style.find(qn("a:lvl%dpPr" % (min(level, 8) + 1)))
    return _rpr_props(lvl.find(qn("a:defRPr"))) if lvl is not None else {}


def _placeholder_defaults(shape, level, layout, master) -> dict:
    """Resolve a placeholder's INHERITED run props at a list level, walking
    layout placeholder -> master placeholder -> master txStyles (most specific
    first). python-pptx doesn't do this; it's why Google/PowerPoint titles come
    back with size=None and render tiny without it."""
    props = {}

    def merge(src):
        for k, v in src.items():
            props.setdefault(k, v)

    try:
        idx = shape.placeholder_format.idx
        pt = str(shape.placeholder_format.type or "")
    except Exception:
        return props
    cat = "title" if "TITLE" in pt else "body"
    try:
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == idx:
                merge(_ph_props(ph, level))
                break
    except Exception:
        pass
    try:
        for ph in master.placeholders:
            mpt = str(ph.placeholder_format.type or "")
            if (cat == "title") == ("TITLE" in mpt):
                merge(_ph_props(ph, level))
                break
    except Exception:
        pass
    merge(_txstyle_props(master, cat, level))
    return props


def _line_dict(shape) -> dict | None:
    """Solid outline color + width (pt), best-effort. None when absent/complex."""
    try:
        ln = shape.line
        color = _rgb_hex(ln.color)
        width = None
        try:
            if ln.width is not None:
                width = round(ln.width / 12700, 2)  # EMU -> pt
        except Exception:
            pass
        if color and (width or 0) > 0:
            return {"color": color, "width": width or 1.0}
    except Exception:
        pass
    return None


# --------------------------------------------------------------------------- #
#  PARSE  pptx -> model                                                       #
# --------------------------------------------------------------------------- #
_ALIGN_TO_STR = {}
_STR_TO_ALIGN = {}


def _init_align_maps():
    from pptx.enum.text import PP_ALIGN

    global _ALIGN_TO_STR, _STR_TO_ALIGN
    _ALIGN_TO_STR = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
    }
    _STR_TO_ALIGN = {v: k for k, v in _ALIGN_TO_STR.items()}


def _run_to_dict(run) -> dict:
    f = run.font
    size = None
    try:
        if f.size is not None:
            size = round(f.size.pt, 1)
    except Exception:
        pass
    return {
        "text": run.text or "",
        # None means "inherit" — filled from the placeholder chain below; the UI
        # and build treat None as falsy/unset, so this stays safe downstream.
        "bold": bool(f.bold) if f.bold is not None else None,
        "italic": bool(f.italic) if f.italic is not None else None,
        "underline": bool(f.underline) if f.underline is not None else False,
        "size": size,
        "color": _rgb_hex(f.color),
        "font": f.name,
    }


def _para_to_dict(para) -> dict:
    runs = [_run_to_dict(r) for r in para.runs]
    if not runs:
        # empty paragraph still needs to occupy a line
        runs = [
            {
                "text": "",
                "bold": False,
                "italic": False,
                "underline": False,
                "size": None,
                "color": None,
                "font": None,
            }
        ]
    align = _ALIGN_TO_STR.get(para.alignment)
    return {"align": align, "level": para.level or 0, "runs": runs}


def _text_frame_to_paras(tf) -> list:
    return [_para_to_dict(p) for p in tf.paragraphs]


def _geom_box(shape, idx) -> dict:
    return {
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "w": emu_to_px(shape.width) or 200.0,
        "h": emu_to_px(shape.height) or 60.0,
        "rot": float(shape.rotation or 0),
        "z": idx,
    }


def _shape_to_element(shape, idx, media_dir, media_rel, slide=None) -> dict | None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    st = shape.shape_type
    base = {"id": _nid("e"), "name": shape.name or ""}
    base.update(_geom_box(shape, idx))

    # ---- picture -------------------------------------------------------- #
    if st == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            ext = (img.ext or "png").lower()
            fname = f"{base['id']}.{ext}"
            with open(os.path.join(media_dir, fname), "wb") as fh:
                fh.write(img.blob)
            base.update({"type": "image", "src": f"{media_rel}/{fname}"})
            return base
        except Exception:
            base.update(
                {
                    "type": "text",
                    "paragraphs": [
                        {
                            "align": "center",
                            "level": 0,
                            "runs": [
                                {
                                    "text": "[image]",
                                    "bold": False,
                                    "italic": False,
                                    "underline": False,
                                    "size": 12,
                                    "color": "#888888",
                                    "font": None,
                                }
                            ],
                        }
                    ],
                }
            )
            return base

    # ---- table ---------------------------------------------------------- #
    if st == MSO_SHAPE_TYPE.TABLE:
        try:
            tbl = shape.table
            rows = []
            for r in tbl.rows:
                rows.append([c.text for c in r.cells])
            base.update({"type": "table", "rows": rows})
            return base
        except Exception:
            return None

    # ---- text / autoshape ---------------------------------------------- #
    if shape.has_text_frame:
        el = dict(base)
        el["type"] = "text"
        el["paragraphs"] = _text_frame_to_paras(shape.text_frame)
        el["fill"] = _fill_hex(shape)
        el["line"] = _line_dict(shape)
        # remember an autoshape geometry name for light fidelity
        try:
            el["geom"] = (
                str(shape.adjustments and shape.auto_shape_type)
                if hasattr(shape, "auto_shape_type")
                else None
            )
        except Exception:
            el["geom"] = None
        # resolve INHERITED run props for placeholders (size/bold/italic/font/color)
        is_ph = False
        try:
            is_ph = bool(shape.is_placeholder)
        except Exception:
            pass
        if is_ph:
            try:
                el["ph"] = str(shape.placeholder_format.type)
            except Exception:
                pass
            cat = "title" if "TITLE" in el.get("ph", "") else "body"
            layout = getattr(slide, "slide_layout", None)
            master = getattr(layout, "slide_master", None) if layout is not None else None
            for para in el["paragraphs"]:
                lvl = para.get("level", 0)
                defs = (
                    _placeholder_defaults(shape, lvl, layout, master) if master is not None else {}
                )
                for r in para.get("runs", []):
                    for k in ("size", "bold", "italic", "font", "color"):
                        if r.get(k) is None and defs.get(k) is not None:
                            r[k] = defs[k]
                    if r.get("size") is None:  # guarantee a size so text is never tiny
                        r["size"] = 40.0 if cat == "title" else 18.0
        # normalize any remaining None bold/italic to False (no inheritance found)
        for para in el["paragraphs"]:
            for r in para.get("runs", []):
                if r.get("bold") is None:
                    r["bold"] = False
                if r.get("italic") is None:
                    r["italic"] = False
        return el

    # ---- fallback: label the unsupported kind so nothing vanishes ------- #
    label = str(st).split()[0].lower() if st is not None else "shape"
    el = dict(base)
    el["type"] = "text"
    el["fill"] = _fill_hex(shape)
    el["paragraphs"] = [
        {
            "align": "center",
            "level": 0,
            "runs": [
                {
                    "text": f"[{label}]",
                    "bold": False,
                    "italic": False,
                    "underline": False,
                    "size": 12,
                    "color": "#9aa0a6",
                    "font": None,
                }
            ],
        }
    ]
    return el


def parse_pptx(path: str, media_dir: str, media_rel: str = "media") -> dict:
    """Parse a .pptx at `path` into the canonical model. Extracts pictures into
    `media_dir` (referenced in the model as `<media_rel>/<file>`)."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "the slides template needs python-pptx: pip install python-pptx"
        ) from None
    _init_align_maps()
    os.makedirs(media_dir, exist_ok=True)
    prs = Presentation(path)

    w_px = emu_to_px(prs.slide_width) or DEFAULT_W_PX
    h_px = emu_to_px(prs.slide_height) or DEFAULT_H_PX

    slides = []
    for s in prs.slides:
        elements = []
        for i, shape in enumerate(s.shapes):
            try:
                el = _shape_to_element(shape, i, media_dir, media_rel, s)
                if el is not None:
                    elements.append(el)
            except Exception as exc:  # never let one bad shape kill the deck
                elements.append(
                    {
                        "id": _nid("e"),
                        "type": "text",
                        "name": "error",
                        "x": 40,
                        "y": 40,
                        "w": 300,
                        "h": 40,
                        "rot": 0,
                        "z": i,
                        "paragraphs": [
                            {
                                "align": "left",
                                "level": 0,
                                "runs": [
                                    {
                                        "text": f"[unreadable shape: {exc}]",
                                        "bold": False,
                                        "italic": False,
                                        "underline": False,
                                        "size": 10,
                                        "color": "#c0392b",
                                        "font": None,
                                    }
                                ],
                            }
                        ],
                    }
                )
        bg = None
        try:
            bg = _fill_hex(s.background) if hasattr(s, "background") else None
        except Exception:
            bg = None
        slides.append({"id": _nid("s"), "background": bg or "#ffffff", "elements": elements})

    name = os.path.splitext(os.path.basename(path))[0]
    return {"schema": 1, "name": name, "width": w_px, "height": h_px, "slides": slides}


# --------------------------------------------------------------------------- #
#  BUILD  model -> pptx                                                        #
# --------------------------------------------------------------------------- #
def _apply_run_font(run, rd):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    f = run.font
    if rd.get("bold") is not None:
        f.bold = bool(rd["bold"])
    if rd.get("italic") is not None:
        f.italic = bool(rd["italic"])
    if rd.get("underline") is not None:
        f.underline = bool(rd["underline"])
    if rd.get("size"):
        try:
            f.size = Pt(float(rd["size"]))
        except Exception:
            pass
    if rd.get("font"):
        f.name = rd["font"]
    col = rd.get("color")
    if col and isinstance(col, str) and col.startswith("#") and len(col) == 7:
        try:
            f.color.rgb = RGBColor.from_string(col[1:].upper())
        except Exception:
            pass


def _write_paragraphs(tf, paragraphs):
    _init_align_maps()
    tf.word_wrap = True
    first = True
    for pd in paragraphs or []:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        al = _STR_TO_ALIGN.get(pd.get("align"))
        if al is not None:
            p.alignment = al
        try:
            p.level = int(pd.get("level") or 0)
        except Exception:
            pass
        for rd in pd.get("runs", []):
            r = p.add_run()
            r.text = rd.get("text", "")
            _apply_run_font(r, rd)


def build_pptx(model: dict, out_path: str, media_root: str) -> str:
    """Rebuild a .pptx from the canonical model. `media_root` is the directory
    that image `src` paths are relative to."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "the slides template needs python-pptx: pip install python-pptx"
        ) from None
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(px_to_emu(model.get("width", DEFAULT_W_PX)))
    prs.slide_height = Emu(px_to_emu(model.get("height", DEFAULT_H_PX)))
    blank = prs.slide_layouts[6]

    for sd in model.get("slides", []):
        slide = prs.slides.add_slide(blank)
        # background
        bg = sd.get("background")
        if bg and isinstance(bg, str) and bg.startswith("#") and len(bg) == 7:
            try:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor.from_string(bg[1:].upper())
            except Exception:
                pass
        els = sorted(sd.get("elements", []), key=lambda e: e.get("z", 0))
        for el in els:
            x, y = Emu(px_to_emu(el.get("x", 0))), Emu(px_to_emu(el.get("y", 0)))
            w, h = Emu(px_to_emu(el.get("w", 100))), Emu(px_to_emu(el.get("h", 40)))
            t = el.get("type")
            if t == "image":
                src = el.get("src") or ""
                p = src if os.path.isabs(src) else os.path.join(media_root, src)
                if os.path.exists(p):
                    try:
                        slide.shapes.add_picture(p, x, y, w, h)
                        continue
                    except Exception:
                        pass
                # fall through to a labelled box if the image is missing
                tb = slide.shapes.add_textbox(x, y, w, h)
                tb.text_frame.text = "[missing image]"
            elif t == "table":
                rows = el.get("rows") or [[""]]
                nr, nc = len(rows), max(len(r) for r in rows)
                gtbl = slide.shapes.add_table(nr, nc, x, y, w, h).table
                for ri, row in enumerate(rows):
                    for ci in range(nc):
                        gtbl.cell(ri, ci).text = row[ci] if ci < len(row) else ""
            else:  # text / shape
                tb = slide.shapes.add_textbox(x, y, w, h)
                fill = el.get("fill")
                if fill and isinstance(fill, str) and fill.startswith("#") and len(fill) == 7:
                    try:
                        tb.fill.solid()
                        tb.fill.fore_color.rgb = RGBColor.from_string(fill[1:].upper())
                    except Exception:
                        pass
                line = el.get("line")
                if line and line.get("color", "").startswith("#"):
                    try:
                        from pptx.util import Pt as _Pt

                        tb.line.color.rgb = RGBColor.from_string(line["color"][1:].upper())
                        tb.line.width = _Pt(float(line.get("width") or 1))
                    except Exception:
                        pass
                tf = tb.text_frame
                va = el.get("valign")
                tf.vertical_anchor = {
                    "top": MSO_ANCHOR.TOP,
                    "middle": MSO_ANCHOR.MIDDLE,
                    "bottom": MSO_ANCHOR.BOTTOM,
                }.get(va, MSO_ANCHOR.TOP)
                _write_paragraphs(tf, el.get("paragraphs"))
                try:
                    if el.get("rot"):
                        tb.rotation = float(el["rot"])
                except Exception:
                    pass
    prs.save(out_path)
    return out_path


# --------------------------------------------------------------------------- #
#  BUILD  model -> self-contained HTML  (no dependency; mirrors the viewer)    #
# --------------------------------------------------------------------------- #
def _esc(s):
    return str(s).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _run_css(r):
    p = []
    if r.get("size"):
        p.append(f"font-size:{r['size'] * 96 / 72:.1f}px")
    if r.get("bold"):
        p.append("font-weight:700")
    if r.get("italic"):
        p.append("font-style:italic")
    if r.get("underline"):
        p.append("text-decoration:underline")
    if r.get("color"):
        p.append(f"color:{r['color']}")
    if r.get("font"):
        p.append(f"font-family:'{r['font']}',sans-serif")
    return ";".join(p)


def build_html(model: dict, media_root: str) -> str:
    """Render the model to a self-contained static HTML deck (images inlined as
    data URIs). No external dependency; slides stack vertically."""
    import base64
    import mimetypes

    W, H = model.get("width", DEFAULT_W_PX), model.get("height", DEFAULT_H_PX)
    fonts = sorted(
        {
            r.get("font")
            for s in model.get("slides", [])
            for e in s.get("elements", [])
            for p in e.get("paragraphs", [])
            for r in p.get("runs", [])
            if r.get("font")
        }
    )
    gf = ""
    if fonts:
        fam = "&".join(
            "family=" + f.replace(" ", "+") + ":ital,wght@0,400;0,700;1,400;1,700" for f in fonts
        )
        gf = f'<link rel="stylesheet" href="https://fonts.googleapis.com/css2?{fam}&display=swap">'

    def img_data(src):
        p = src if os.path.isabs(src) else os.path.join(media_root, src)
        try:
            with open(p, "rb") as fh:
                b = base64.b64encode(fh.read()).decode()
            mt = mimetypes.guess_type(p)[0] or "image/png"
            return f"data:{mt};base64,{b}"
        except Exception:
            return ""

    slides_html = []
    for s in model.get("slides", []):
        els = []
        for el in sorted(s.get("elements", []), key=lambda e: e.get("z", 0)):
            box = (
                f"position:absolute;left:{el.get('x', 0)}px;top:{el.get('y', 0)}px;"
                f"width:{el.get('w', 0)}px;height:{el.get('h', 0)}px;overflow:hidden"
            )
            if el.get("rot"):
                box += f";transform:rotate({el['rot']}deg)"
            if el["type"] == "image":
                els.append(
                    f'<div style="{box}"><img src="{img_data(el.get("src", ""))}" '
                    f'style="width:100%;height:100%;object-fit:contain"></div>'
                )
            elif el["type"] == "table":
                rows = "".join(
                    "<tr>" + "".join(f"<td>{_esc(c)}</td>" for c in row) + "</tr>"
                    for row in el.get("rows", [])
                )
                els.append(f'<div style="{box}"><table class="t">{rows}</table></div>')
            else:
                va = {"middle": "center", "bottom": "flex-end"}.get(el.get("valign"), "flex-start")
                fill = f";background:{el['fill']}" if el.get("fill") else ""
                line = (
                    f";border:{el['line'].get('width', 1)}px solid {el['line']['color']}"
                    if el.get("line") and el["line"].get("color")
                    else ""
                )
                paras = ""
                for pa in el.get("paragraphs", []):
                    runs = "".join(
                        f'<span style="{_run_css(r)}">{_esc(r.get("text", "")) or "&#8203;"}</span>'
                        for r in pa.get("runs", [])
                    )
                    paras += f'<div style="text-align:{pa.get("align") or "left"}">{runs}</div>'
                els.append(
                    f'<div style="{box};display:flex;flex-direction:column;'
                    f'justify-content:{va};padding:2px 4px{fill}{line}">{paras}</div>'
                )
        bg = s.get("background") or "#ffffff"
        slides_html.append(
            f'<div class="page"><section class="slide" style="width:{W}px;'
            f'height:{H}px;background:{bg}">{"".join(els)}</section></div>'
        )
    return f"""<!doctype html><html><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{_esc(model.get("name", "Deck"))}</title>{gf}
<style>
  :root{{--W:{W}px;--H:{H}px}}
  body{{margin:0;background:#f4f4f2;font-family:system-ui,-apple-system,Segoe UI,Roboto,sans-serif;
    display:flex;flex-direction:column;align-items:center;gap:22px;padding:22px}}
  .page{{width:min(1280px,96vw);position:relative}}
  .slide{{position:relative;overflow:hidden;transform-origin:top left;
    box-shadow:0 1px 3px rgba(0,0,0,.1),0 6px 20px rgba(0,0,0,.06)}}
  .slide .t{{border-collapse:collapse;width:100%;height:100%;font-size:13px}}
  .slide .t td{{border:1px solid #cfd3d8;padding:3px 6px}}
  #hint{{position:fixed;bottom:12px;right:14px;font-size:12px;color:#9a9a9a}}
  body.present{{background:#000;padding:0;gap:0;justify-content:center}}
  body.present .page{{display:none}} body.present .page.show{{display:grid;place-items:center;
    width:100vw;height:100vh}}
  body.present #hint{{color:#666}}
</style></head><body>{"".join(slides_html)}
<div id="hint">Press F to present · ←/→ to navigate</div>
<script>
  var W={W},H={H},pages=[].slice.call(document.querySelectorAll('.page')),i=0,present=false;
  function fit(){{pages.forEach(function(pg){{var s=pg.querySelector('.slide');
    var w=present?innerWidth:pg.clientWidth, sc=Math.min(w/W, (present?innerHeight:1e9)/H);
    s.style.transform='scale('+sc+')'; pg.style.height=(H*sc)+'px';}});}}
  function show(n){{i=Math.max(0,Math.min(n,pages.length-1));
    pages.forEach(function(p,k){{p.classList.toggle('show',k===i);}}); fit();}}
  function enter(){{present=true;document.body.classList.add('present');show(i);
    if(document.documentElement.requestFullscreen)document.documentElement.requestFullscreen();}}
  function exit(){{present=false;document.body.classList.remove('present');
    pages.forEach(function(p){{p.classList.remove('show');}});fit();
    if(document.fullscreenElement)document.exitFullscreen();}}
  addEventListener('resize',fit); fit();
  addEventListener('keydown',function(e){{
    if(e.key==='f'||e.key==='F'){{present?exit():enter();}}
    else if(present&&(e.key==='ArrowRight'||e.key===' ')){{e.preventDefault();show(i+1);}}
    else if(present&&e.key==='ArrowLeft'){{show(i-1);}}
    else if(present&&e.key==='Escape'){{exit();}}
  }});
  addEventListener('click',function(){{if(present)show(i+1);}});
</script></body></html>"""


# --------------------------------------------------------------------------- #
#  BUILD  model -> PDF  (fpdf2; draws the canonical model, sibling of build_*) #
# --------------------------------------------------------------------------- #
_PX2PT = 72 / 96
_SERIF_HINTS = (
    "times",
    "georgia",
    "serif",
    "garamond",
    "playfair",
    "merriweather",
    "lora",
    "noto serif",
    "roboto slab",
    "source serif",
    "spectral",
    "cormorant",
    "baskerville",
    "bodoni",
    "domine",
    "bitter",
    "alegreya",
)
_MONO_HINTS = ("mono", "courier", "code", "consol")


def _hex_rgb(h, default=(32, 32, 32)):
    try:
        h = h.lstrip("#")
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return default


def _run_core_font(r):
    name = (r.get("font") or "").lower()
    if any(k in name for k in _MONO_HINTS):
        fam = "courier"
    elif any(k in name for k in _SERIF_HINTS):
        fam = "times"
    else:
        fam = "helvetica"
    style = ("B" if r.get("bold") else "") + ("I" if r.get("italic") else "")
    return fam, style


def _emit_line(pdf, line, left, maxw, align, cy):
    """Draw one wrapped line of (word, fam, style, size, rgb); return the new y."""

    def wd(w):
        pw, fam, st, sz, col = w
        pdf.set_font(fam, st, sz)
        return pdf.get_string_width(pw)

    space = 0
    pdf.set_font("helvetica", "", max(w[3] for w in line))
    sp = pdf.get_string_width(" ")
    linew = sum(wd(w) for w in line) + sp * (len(line) - 1)
    if align == "center":
        bx = left + (maxw - linew) / 2
    elif align == "right":
        bx = left + (maxw - linew)
    else:
        bx = left
    baseline = cy + max(w[3] for w in line) * 0.92
    for idx, w in enumerate(line):
        pw, fam, st, sz, col = w
        pdf.set_font(fam, st, sz)
        pdf.set_text_color(*col)
        if idx:
            bx += sp
        pdf.text(bx, baseline, pw)
        bx += pdf.get_string_width(pw)
    return cy + max(w[3] for w in line) * 1.25


def _draw_text_box(pdf, el, x, y, w, h):
    pad = 3
    left, top, maxw, bottom = x + pad, y + pad, max(4, w - 2 * pad), y + h
    cy = top
    for para in el.get("paragraphs", []):
        align = para.get("align") or "left"
        words = []
        for r in para.get("runs", []):
            fam, st = _run_core_font(r)
            sz = float(r.get("size") or 18)
            col = _hex_rgb(r.get("color")) if r.get("color") else (32, 32, 32)
            for pw in (r.get("text") or "").split(" "):
                words.append((pw, fam, st, sz, col))
        if not words:
            cy += 14
            continue
        line, linew = [], 0
        pdf.set_font("helvetica", "", 12)
        for wt in words:
            pdf.set_font(wt[1], wt[2], wt[3])
            ww = pdf.get_string_width(wt[0])
            sp = pdf.get_string_width(" ") if line else 0
            if line and linew + sp + ww > maxw:
                cy = _emit_line(pdf, line, left, maxw, align, cy)
                line, linew = [], 0
                sp = 0
            line.append(wt)
            linew += sp + ww
            if cy > bottom:
                return
        if line:
            cy = _emit_line(pdf, line, left, maxw, align, cy)


def build_pdf(model: dict, out_path: str, media_root: str) -> str:
    try:
        from fpdf import FPDF
    except ImportError:
        raise ImportError("PDF export needs fpdf2: pip install fpdf2") from None
    W = model.get("width", DEFAULT_W_PX) * _PX2PT
    H = model.get("height", DEFAULT_H_PX) * _PX2PT
    pdf = FPDF(unit="pt", format=(W, H))
    pdf.set_auto_page_break(False)
    for s in model.get("slides", []):
        pdf.add_page()
        bg = s.get("background")
        if bg:
            pdf.set_fill_color(*_hex_rgb(bg, (255, 255, 255)))
            pdf.rect(0, 0, W, H, style="F")
        for el in sorted(s.get("elements", []), key=lambda e: e.get("z", 0)):
            x, y = el.get("x", 0) * _PX2PT, el.get("y", 0) * _PX2PT
            w, h = el.get("w", 0) * _PX2PT, el.get("h", 0) * _PX2PT
            t = el.get("type")
            if t == "image":
                src = el.get("src") or ""
                p = src if os.path.isabs(src) else os.path.join(media_root, src)
                if os.path.exists(p):
                    try:
                        pdf.image(p, x, y, w, h)
                    except Exception:
                        pass
            elif t == "table":
                rows = el.get("rows") or []
                nr = len(rows) or 1
                rh = h / nr
                for ri, row in enumerate(rows):
                    nc = len(row) or 1
                    cw = w / nc
                    for ci, cell in enumerate(row):
                        pdf.set_draw_color(200, 205, 210)
                        pdf.rect(x + ci * cw, y + ri * rh, cw, rh)
                        _draw_text_box(
                            pdf,
                            {
                                "paragraphs": [
                                    {"align": "left", "runs": [{"text": cell, "size": 11}]}
                                ]
                            },
                            x + ci * cw,
                            y + ri * rh,
                            cw,
                            rh,
                        )
            else:
                if el.get("fill"):
                    pdf.set_fill_color(*_hex_rgb(el["fill"], (255, 255, 255)))
                    pdf.rect(x, y, w, h, style="F")
                _draw_text_box(pdf, el, x, y, w, h)
    pdf.output(out_path)
    return out_path


# --------------------------------------------------------------------------- #
#  BUILD  model -> Markdown outline                                           #
# --------------------------------------------------------------------------- #
def build_md(model: dict) -> str:
    out = [f"# {model.get('name', 'Deck')}", ""]
    for i, s in enumerate(model.get("slides", [])):
        out.append(f"## Slide {i + 1}")
        texts = [
            e
            for e in sorted(s.get("elements", []), key=lambda e: e.get("y", 0))
            if e.get("type") == "text"
        ]
        for j, el in enumerate(texts):
            lines = [
                "".join(r.get("text", "") for r in p.get("runs", []))
                for p in el.get("paragraphs", [])
            ]
            lines = [ln for ln in lines if ln.strip()]
            if not lines:
                continue
            if j == 0:
                out.append(f"### {' '.join(lines)}")
            else:
                out.extend(f"- {ln}" for ln in lines)
        for el in s.get("elements", []):
            if el.get("type") == "table":
                for row in el.get("rows", []):
                    out.append("| " + " | ".join(row) + " |")
            elif el.get("type") == "image":
                out.append(f"![image]({el.get('src', '')})")
        out.append("")
    return "\n".join(out)


# --------------------------------------------------------------------------- #
#  MODEL OPS  (shared canonical mutations — used by the AI actions and the UI) #
# --------------------------------------------------------------------------- #
def find_slide(model, slide_id):
    for s in model.get("slides", []):
        if s["id"] == slide_id:
            return s
    return None


def find_element(model, el_id):
    for s in model.get("slides", []):
        for e in s.get("elements", []):
            if e["id"] == el_id:
                return s, e
    return None, None


def _next_z(slide):
    return (max([e.get("z", 0) for e in slide.get("elements", [])], default=-1)) + 1


def new_text_element(
    x=100, y=100, w=400, h=80, text="Text", size=18, align="left", color="#202124", bold=False
):
    return {
        "id": _nid("e"),
        "type": "text",
        "name": "TextBox",
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "rot": 0,
        "z": 0,
        "fill": None,
        "valign": "top",
        "paragraphs": [
            {
                "align": align,
                "level": 0,
                "runs": [
                    {
                        "text": text,
                        "bold": bold,
                        "italic": False,
                        "underline": False,
                        "size": size,
                        "color": color,
                        "font": None,
                    }
                ],
            }
        ],
    }


def new_image_element(src, x=100, y=100, w=400, h=300):
    return {
        "id": _nid("e"),
        "type": "image",
        "name": "Picture",
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "rot": 0,
        "z": 0,
        "src": src,
    }


def new_slide(background="#ffffff"):
    return {"id": _nid("s"), "background": background, "elements": []}


def blank_model(name="Untitled", w=DEFAULT_W_PX, h=DEFAULT_H_PX):
    return {"schema": 1, "name": name, "width": w, "height": h, "slides": [new_slide()]}


# --------------------------------------------------------------------------- #
#  CLI smoke test                                                              #
# --------------------------------------------------------------------------- #
if __name__ == "__main__":
    import json
    import sys

    if len(sys.argv) < 2:
        print("usage: engine.py <file.pptx> [out.pptx]")
        raise SystemExit(1)
    tmp = os.path.join(os.path.dirname(os.path.abspath(sys.argv[1])), "_engine_media")
    m = parse_pptx(sys.argv[1], tmp)
    print(
        json.dumps(
            {
                "name": m["name"],
                "size": [m["width"], m["height"]],
                "n_slides": len(m["slides"]),
                "el_types": [[e["type"] for e in s["elements"]] for s in m["slides"]],
            },
            indent=2,
        )
    )
    if len(sys.argv) > 2:
        build_pptx(m, sys.argv[2], tmp)
        print("rebuilt ->", sys.argv[2])
