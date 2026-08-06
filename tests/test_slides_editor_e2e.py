"""End-to-end tests for the slides editor's text-format toolbar and selection,
driven through Playwright against a real server (FUSED_RENDER_CORE_TEMPLATES
points at this checkout's templates). These pin the behaviours reported broken
on the create-new-flows branch:

  * picking a Style preset applies its size to the selected text box,
  * picking a Font applies it,
  * the Bold button toggles,
  * Ctrl+A with a text box selected selects that box's TEXT (enters inline edit)
    instead of selecting every element on the slide.

Skipped when playwright, python-pptx, or the built React shell is missing.

Run serially: PYTHONPATH=<checkout> python -m pytest tests/test_slides_editor_e2e.py -o addopts=""
"""
import os
import socket
import subprocess
import sys
import time
import urllib.request
from urllib.parse import quote

import pytest

pytest.importorskip("playwright.sync_api")
pytest.importorskip("pptx")
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHELL = os.path.join(ROOT, "fused_render", "static", "shell-dist", "index.html")

pytestmark = pytest.mark.skipif(
    not os.path.exists(SHELL), reason="React shell not built")


def _free_port():
    with socket.socket() as s:
        s.bind(("127.0.0.1", 0))
        return s.getsockname()[1]


def _sample_pptx(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tf = tb.text_frame
    tf.text = "Hello world"
    tf.add_paragraph().text = "Second line"
    prs.save(path)


@pytest.fixture(scope="module")
def server(tmp_path_factory):
    home = tmp_path_factory.mktemp("slides-home")
    work = tmp_path_factory.mktemp("slides-work")
    port = _free_port()
    env = dict(os.environ)
    env["PYTHONPATH"] = ROOT
    env["FUSED_RENDER_HOME"] = str(home)
    env["FUSED_RENDER_CORE_TEMPLATES"] = os.path.join(ROOT, "fused_render", "templates")
    log_path = home / "server.log"
    with open(log_path, "ab") as log:
        proc = subprocess.Popen(
            [sys.executable, "-m", "fused_render.cli", "serve",
             "--port", str(port), "--no-browser", "--start-dir", str(work)],
            env=env, stdout=log, stderr=subprocess.STDOUT)
    deadline = time.time() + 60
    while time.time() < deadline:
        try:
            urllib.request.urlopen(f"http://127.0.0.1:{port}/api/config", timeout=2).read()
            break
        except OSError:
            if proc.poll() is not None:
                raise RuntimeError(
                    "server exited:\n" + log_path.read_text("utf-8", errors="replace")[-2000:])
            time.sleep(0.3)
    else:
        proc.kill()
        raise RuntimeError("server did not come up")
    yield {"port": port, "work": work}
    proc.kill()
    proc.wait(timeout=15)


def _wait_for(fn, timeout, step=0.2):
    deadline = time.time() + timeout
    while time.time() < deadline:
        value = fn()
        if value:
            return value
        time.sleep(step)
    raise AssertionError(f"condition not met within {timeout}s")


def _find_frame(page):
    def frame_or_escape():
        page.keyboard.press("Escape")  # dismiss the shell welcome tour
        return next((f for f in page.frames if "/render" in f.url), None)
    return _wait_for(frame_or_escape, 30, step=0.5)


# Synthetic-event helpers injected into the frame: the canvas elements are
# CSS-transform-scaled and the pickers use mousedown handlers, so we drive them
# the way the app itself dispatches, then read state back.
_HELPERS = r"""
window.__t = (() => {
  const q=(s)=>document.querySelector(s), qa=(s)=>[...document.querySelectorAll(s)];
  const w=window;
  const fire=(n,t)=>{const r=n.getBoundingClientRect();
    n.dispatchEvent(new MouseEvent(t,{bubbles:true,cancelable:true,
      clientX:r.left+r.width/2,clientY:r.top+r.height/2,view:w}));};
  const md=(n)=>n.dispatchEvent(new MouseEvent('mousedown',{bubbles:true,cancelable:true,view:w}));
  const clk=(n)=>n.dispatchEvent(new MouseEvent('click',{bubbles:true,cancelable:true,view:w}));
  const firstText=()=>qa('#canvas .el-text')[0];
  // a picker pop must sit just under its trigger, left-aligned — not drop to the
  // bottom of the slide (which happens when its container isn't position:relative)
  const geom=(trig,pop)=>{const t=trig.getBoundingClientRect(), p=pop.getBoundingClientRect();
    return {anchoredBelow: p.top>=t.bottom-1 && p.top<=t.bottom+8,
            leftAligned: Math.abs(p.left-t.left)<3,
            droppedAway: p.top>t.bottom+100}; };
  return {
    selectBox(){ const e=firstText(); fire(e,'mousedown'); fire(e,'mouseup');
      return q('#canvas').querySelectorAll('.el.selected').length; },
    pickStyle(label){ clk(q('#styleTrigger'));
      const o=qa('#styleList > *').find(x=>x.textContent.trim().replace('✓','')===label);
      if(o) md(o); return q('#fontSize').value; },
    pickFont(name){ clk(q('#fontTrigger'));
      const o=qa('#fontList .pop-item').find(x=>x.textContent.includes(name));
      if(o) md(o); return q('#fontName').textContent; },
    boldOn(){ return q('#bBtn').classList.contains('on'); },
    clickBold(){ clk(q('#bBtn')); return q('#bBtn').classList.contains('on'); },
    ctrlA(){ document.dispatchEvent(new KeyboardEvent('keydown',
      {key:'a',code:'KeyA',ctrlKey:true,bubbles:true,cancelable:true})); },
    styleGeom(){ clk(q('#styleTrigger')); return geom(q('#styleTrigger'), q('#stylePop')); },
    fontGeom(){ clk(q('#fontTrigger')); return geom(q('#fontTrigger'), q('#fontPop')); },
    enterEdit(){ const e=firstText(); fire(e,'mousedown'); fire(e,'mouseup'); fire(e,'dblclick');
      const ed=qa('#canvas .el-text[contenteditable="true"]')[0]; if(ed) ed.focus();
      return qa('#canvas .el-text[contenteditable="true"]').length; },
    editableCount(){ return qa('#canvas .el-text[contenteditable="true"]').length; },
    selectedEls(){ return q('#canvas').querySelectorAll('.el.selected').length; },
    clickEmpty(){ const r=q('#canvas').getBoundingClientRect(), x=r.left+8, y=r.top+r.height-8;
      const cv=q('#canvas');
      cv.dispatchEvent(new MouseEvent('mousedown',{bubbles:true,cancelable:true,view:w,clientX:x,clientY:y}));
      cv.dispatchEvent(new MouseEvent('mouseup',{bubbles:true,cancelable:true,view:w,clientX:x,clientY:y}));
      return qa('#canvas .el-text[contenteditable="true"]').length; },
    pressEscape(){ document.dispatchEvent(new KeyboardEvent('keydown',{key:'Escape',bubbles:true,cancelable:true})); },
    dragBox(dx,dy){ const e=firstText(); const b=e.getBoundingClientRect();
      const cx=b.left+b.width/2, cy=b.top+b.height/2;
      e.dispatchEvent(new MouseEvent('mousedown',{bubbles:true,cancelable:true,view:w,clientX:cx,clientY:cy}));
      document.dispatchEvent(new MouseEvent('mousemove',{bubbles:true,view:w,clientX:cx+dx,clientY:cy+dy}));
      document.dispatchEvent(new MouseEvent('mouseup',{bubbles:true,view:w,clientX:cx+dx,clientY:cy+dy}));
      return Math.round(firstText().getBoundingClientRect().left-b.left); },
    editingCount(){ return qa('#canvas .el-text[contenteditable="true"]').length; },
    selectedEls(){ return q('#canvas').querySelectorAll('.el.selected').length; },
    textSelLen(){ return (w.getSelection()+'').length; },
    fontSizeVal(){ return q('#fontSize').value; },
  };
})();
"""


@pytest.fixture(scope="module")
def frame(server):
    pptx_path = os.path.join(str(server["work"]), "deck.pptx")
    _sample_pptx(pptx_path)
    url_path = "/".join(quote(seg, safe="") for seg in
                        pptx_path.replace("\\", "/").split("/") if seg)
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        page = browser.new_page()
        page.goto(f"http://127.0.0.1:{server['port']}/view/{url_path}?mode=edit")
        fr = _find_frame(page)
        _wait_for(lambda: fr.locator("#canvas .el-text").count() or None, 30)
        fr.evaluate(_HELPERS)
        yield fr
        browser.close()


def test_style_preset_applies_size(frame):
    assert frame.evaluate("window.__t.selectBox()") == 1
    # Caption preset is 12pt; the size box reflects the applied run size.
    assert frame.evaluate("window.__t.pickStyle('Caption')") == "12"
    frame.evaluate("window.__t.selectBox()")
    assert frame.evaluate("window.__t.pickStyle('Title')") == "44"


def test_font_applies(frame):
    frame.evaluate("window.__t.selectBox()")
    assert frame.evaluate("window.__t.pickFont('Georgia')") == "Georgia"


def test_pickers_anchor_under_trigger(frame):
    # both pops must open directly under their trigger, not drop to the bottom
    frame.evaluate("window.__t.selectBox()")
    s = frame.evaluate("window.__t.styleGeom()")
    assert s["anchoredBelow"] and s["leftAligned"] and not s["droppedAway"], s
    f = frame.evaluate("window.__t.fontGeom()")
    assert f["anchoredBelow"] and f["leftAligned"] and not f["droppedAway"], f


def test_bold_toggles(frame):
    frame.evaluate("window.__t.selectBox()")
    before = frame.evaluate("window.__t.boldOn()")
    after = frame.evaluate("window.__t.clickBold()")
    assert after != before


def test_click_out_of_edit_leaves_box_movable(frame):
    # edit a box, then click empty canvas: the edit must commit (not stay stuck
    # in contenteditable), and the box must be re-selectable and draggable.
    assert frame.evaluate("window.__t.enterEdit()") == 1
    assert frame.evaluate("window.__t.clickEmpty()") == 0        # committed, no editable left
    assert frame.evaluate("window.__t.selectBox()") == 1         # re-selectable
    assert frame.evaluate("window.__t.dragBox(40, 0)") >= 25     # and draggable


def test_escape_exits_edit_and_keeps_box_selected(frame):
    assert frame.evaluate("window.__t.enterEdit()") == 1
    frame.evaluate("window.__t.pressEscape()")
    assert _wait_for(lambda: frame.evaluate("window.__t.editableCount()") == 0 or None, 3)
    assert frame.evaluate("window.__t.selectedEls()") == 1       # stays selected
    assert frame.evaluate("window.__t.dragBox(40, 0)") >= 25     # movable right away


def test_ctrl_a_selects_box_text_not_all_elements(frame):
    # a single text box selected, not editing
    assert frame.evaluate("window.__t.selectBox()") == 1
    frame.evaluate("window.__t.ctrlA()")
    # Ctrl+A enters inline edit on that box and selects its text …
    assert _wait_for(lambda: frame.evaluate("window.__t.editingCount()") == 1 or None, 3)
    assert frame.evaluate("window.__t.textSelLen()") > 0
    # … and does NOT turn into a select-all-elements (only the one box stays selected)
    assert frame.evaluate("window.__t.selectedEls()") == 1
